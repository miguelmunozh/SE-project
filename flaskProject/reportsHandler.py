import os
from datetime import date

import xlsxwriter
from docx import Document
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from log.logHandler import LogHandler
from analyst.analystHandler import AnalystHandler

logsHandler = LogHandler()
analystHandler = AnalystHandler()

def generateERB(event, findingsList, systemsList):
    # Code to generate ERB Report
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    left = top = Inches(.5)
    slide.shapes.add_picture('images/Capture.PNG', left, top)

    title.text = "U.S ARMY COMBAT CAPABILITIES DEVELOPMENT COMMAND - DATA & ANALYSIS CENTER"
    subtitle.text = event.getName()
    title.text_frame.paragraphs[0].font.size = Pt(30)

    # a second page
    title_slide_layout1 = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout1)
    body_shape = slide.shapes.placeholders[1]

    left = top = Inches(.5)
    height = Inches(1)
    slide.shapes.add_picture('images/Capture1.PNG', left, top, height=height)
    left = Inches(7.4)
    slide.shapes.add_picture('images/Capture2.PNG', left, top, height=height)
    title1 = slide.shapes.title
    title1.text = "SCOPE"
    title1.text_frame.paragraphs[0].font.size = Pt(30)

    tf = body_shape.text_frame
    tf.text = 'Systems assessed during the CVPA are as follows:'

    # get just the non-archived systems
    for system in systemsList:
        p = tf.add_paragraph()
        p.text = system.getName()
        p.level = 1

    # third page
    title_slide_layout1 = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout1)
    body_shape1 = slide.shapes.placeholders[0]

    left = top = Inches(.5)
    height = Inches(1)
    slide.shapes.add_picture('images/Capture1.PNG', left, top, height=height)
    left = Inches(7.4)
    slide.shapes.add_picture('images/Capture2.PNG', left, top, height=height)
    title1 = slide.shapes.title
    title1.text = "FINDINGS"
    title1.text_frame.paragraphs[0].font.size = Pt(30)

    # create table and then add a row per each finding
    cols = 5
    rows = 10
    top = Inches(1.7)
    left = Inches(0.5)
    width = Inches(9.0)
    height = Inches(1)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(1.8)
    table.columns[4].width = Inches(1.8)

    # write column headings
    table.cell(0, 0).text = 'ID'
    table.cell(0, 1).text = 'System'
    table.cell(0, 2).text = 'Finding'
    table.cell(0, 3).text = 'Impact'
    table.cell(0, 4).text = 'Risk'

    row = 1
    for finding in findingsList:
        # write body cells
        table.cell(row, 0).text = str(finding.getid())
        table.cell(row, 1).text = finding.getHostName()
        table.cell(row, 2).text = finding.getType()
        table.cell(row, 3).text = str(finding.getImpactLevel())
        table.cell(row, 4).text = finding.getRisk()
        row += 1

    # change font size for the table
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)

    # fourth page
    for finding in findingsList:

        title_slide_layout1 = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_slide_layout1)
        body_shape1 = slide.shapes.placeholders[1]
        left = top = Inches(.5)
        height = Inches(1)
        slide.shapes.add_picture('images/Capture1.PNG', left, top, height=height)
        left = Inches(7.4)
        slide.shapes.add_picture('images/Capture2.PNG', left, top, height=height)
        title1 = slide.shapes.title
        title1.text = finding.getType()
        title1.text_frame.paragraphs[0].font.size = Pt(30)

        # create table and then add a row per each finding
        cols = 10
        rows = 14
        top = Inches(1.5)
        left = Inches(0.5)
        width = Inches(11.0)
        height = Inches(1)

        fTable = slide.shapes.add_table(rows, cols, left, top, width, height).table
        fTable.columns[0].width = Inches(.5)
        fTable.columns[3].width = Inches(1.7)
        fTable.columns[4].width = Inches(1)
        fTable.columns[7].width = Inches(.5)
        fTable.columns[8].width = Inches(.5)
        fTable.columns[9].width = Inches(.5)

        fTable.cell(1, 0).text = 'ID'
        fTable.cell(1, 1).text = str(finding.getid())

        fTable.cell(2, 0).text = 'Host Names'
        fTable.cell(3, 0).text = str(finding.getHostName())
        fTable.cell(2, 2).text = 'IP:PORT'
        fTable.cell(3, 2).text = str(finding.getIpPort())

        fTable.cell(1, 3).text = 'Impact Score'
        fTable.cell(1, 4).text = str(finding.getImpactScore())
        fTable.cell(2, 3).text = 'Cat'
        fTable.cell(2, 4).text = str(finding.getSeverityCategoryCode().name)
        fTable.cell(3, 3).text = 'Cat Score'
        fTable.cell(3, 4).text = str(finding.getSeverityCategoryScore())
        fTable.cell(4, 3).text = 'Vs Score'
        fTable.cell(4, 4).text = str(finding.getVulnerabilitySeverity())
        fTable.cell(5, 3).text = 'Vq'
        fTable.cell(5, 4).text = str(finding.getQualitativeVulnerabilitySeverity())
        fTable.cell(6, 3).text = 'Impact Rationale'
        fTable.cell(6, 4).text = str(finding.getImpactDescription())

        fTable.cell(1, 5).text = 'Status'
        fTable.cell(1, 6).text = finding.getStatus().name
        fTable.cell(2, 5).text = 'Likelihood'
        fTable.cell(2, 6).text = finding.getLikelihood()
        fTable.cell(3, 5).text = 'Impact'
        fTable.cell(3, 6).text = finding.getImpactLevel().name
        fTable.cell(4, 5).text = 'Risk'
        fTable.cell(4, 6).text = finding.getRisk()
        fTable.cell(5, 5).text = 'CM'
        fTable.cell(5, 6).text = str(finding.getCountermeasureEffectivenessRating().name)

        fTable.cell(1, 7).text = 'Posture'
        fTable.cell(2, 7).text = finding.getPosture()

        fTable.cell(3, 7).text = 'C'
        fTable.cell(4, 7).text = finding.getConfidentiality()
        fTable.cell(3, 8).text = 'I'
        fTable.cell(4, 8).text = finding.getIntegrity()
        fTable.cell(3, 9).text = 'A'
        fTable.cell(4, 9).text = finding.getAvailability()

        fTable.cell(7, 1).text = 'Type'
        fTable.cell(7, 2).text = finding.getType()

        fTable.cell(8, 2).text = finding.getDescription()
        fTable.cell(9, 0).text = 'Description'
        fTable.cell(9, 1).text = finding.getLongDescription()

        fTable.cell(10, 2).text = finding.getMitigationBriefDescription()
        fTable.cell(11, 0).text = 'Mitigation'
        fTable.cell(11, 1).text = finding.getMitigationLongDescription()

        fTable.cell(12, 0).text = 'Reference'
        fTable.cell(12, 2).text = 'Figure x'
        fTable.cell(13, 0).text = 'C-CONFIDENTIALITY  I-INTEGRITY  A-AVAILABILITY  CM-COUNTERMEASURE'

        # expand the column from 0 to 1 at row 2
        cell = fTable.cell(2, 0)
        other_cell = fTable.cell(2, 1)
        cell.merge(other_cell)
        # expand the column from 0 to 1
        cell = fTable.cell(3, 0)
        other_cell = fTable.cell(6, 1)
        cell.merge(other_cell)

        cell = fTable.cell(3, 2)
        other_cell = fTable.cell(6, 2)
        cell.merge(other_cell)

        cell = fTable.cell(6, 4)
        other_cell = fTable.cell(6, 6)
        cell.merge(other_cell)

        cell = fTable.cell(1, 7)
        other_cell = fTable.cell(1, 9)
        cell.merge(other_cell)
        cell = fTable.cell(2, 7)
        other_cell = fTable.cell(2, 9)
        cell.merge(other_cell)

        cell = fTable.cell(7, 0)
        other_cell = fTable.cell(7, 1)
        cell.merge(other_cell)

        cell = fTable.cell(7, 2)
        other_cell = fTable.cell(7, 9)
        cell.merge(other_cell)

        cell = fTable.cell(8, 1)
        other_cell = fTable.cell(8, 9)
        cell.merge(other_cell)
        cell = fTable.cell(10, 1)
        other_cell = fTable.cell(10, 9)
        cell.merge(other_cell)

        cell = fTable.cell(9, 1)
        other_cell = fTable.cell(9, 9)
        cell.merge(other_cell)
        cell = fTable.cell(11, 1)
        other_cell = fTable.cell(11, 9)
        cell.merge(other_cell)

        cell = fTable.cell(12, 0)
        other_cell = fTable.cell(12, 1)
        cell.merge(other_cell)

        cell = fTable.cell(12, 2)
        other_cell = fTable.cell(12, 9)
        cell.merge(other_cell)

        cell = fTable.cell(13, 0)
        other_cell = fTable.cell(13, 9)
        cell.merge(other_cell)

        for cell in iter_cells(fTable):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

    # histogram page
    title_slide_layout1 = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout1)
    body_shape = slide.shapes.placeholders[1]

    left = top = Inches(.5)
    height = Inches(1)
    slide.shapes.add_picture('images/Capture1.PNG', left, top, height=height)
    left = Inches(7.6)
    slide.shapes.add_picture('images/Capture2.PNG', left, top, height=height)
    title1 = slide.shapes.title
    title1.text = "Findings Histogram"
    title1.text_frame.paragraphs[0].font.size = Pt(30)

    subtitle = slide.placeholders[1]
    subtitle.text = "Finding Risk"

    # define chart data ---------------------
    info = 0
    veryLow = 0
    low = 0
    medium = 0
    high = 0
    veryHigh = 0
    for f in findingsList:
        if f.getRisk() == 'I':
            info += 1
        if f.getRisk() == 'VL':
            veryLow += 1
        if f.getRisk() == 'L':
            low += 1
        if f.getRisk() == 'M':
            medium += 1
        if f.getRisk() == 'H':
            high += 1
        if f.getRisk() == 'VH':
            veryHigh += 1

    chart_data = CategoryChartData()
    chart_data.categories = ['INFO', 'VERY LOW', 'LOW', 'MEDIUM', 'HIGH', 'VERY HIGH']
    chart_data.add_series('Series 1', (info, veryLow, low, medium, high, veryHigh))

    # add chart to slide --------------------
    x, y, cx, cy = Inches(1), Inches(2.5), Inches(8), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )

    prs.save(os.path.expanduser('~/Downloads/ERBReport.pptx'))


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def generateFinalTecReport(event, findingsList):
    # Code to generate RiskMatrixReport
    document = Document()
    document.add_picture('images/Capture4.PNG', width=Inches(4))
    today = date.today()
    section = document.sections[0]
    header = section.header
    paragraph = header.paragraphs[0]
    paragraph.text = "\t\t" + today.strftime("%B %d, %Y")
    document.add_heading("Combat Capabilities Development Command (CCDC) Data &amp; Analysis Center (DAC) Enter "
                         "System Name Enter Event Type (e.g., CVPA, CVI, VoF, etc) Report", level=1)

    # get event team, get the analysts ann their first and last name here
    names = []
    for analyst in analystHandler.getAllAnalyst():
        for initial in event.getEventTeam():
            if initial == analyst.getInitial():
                mystr = analyst.getFirstName()
                mystr += " "
                mystr += analyst.getLastName()
                names.append(mystr)
    my_string = ', '.join(names)
    p = document.add_paragraph()
    p.add_run("by " + my_string).bold = True

    document.add_paragraph("classified by:" + event.getClassifiedBy())
    document.add_paragraph("Derived from:" + event.getSecurityClassificationTitleGuide())
    document.add_paragraph("Declassify on:" + event.getDeclassificationDate())

    # add a new page
    document.add_page_break()
    document.add_heading("DESTRUCTION NOTICE")
    document.add_paragraph("Destroy by any method that will prevent disclosure of contents or reconstruction of the "
                           "document.")
    document.add_heading("DISCLAIMER")
    document.add_paragraph(
        "The findings in this report are not to be construed as an official Department of the Army position unless so specified by other official documentation.")
    document.add_heading("WARNING")
    document.add_paragraph(
        "Information and data contained in this document are based on the input available at the time of preparation.")
    document.add_heading("TRADE NAMES")
    document.add_paragraph(
        "The use of trade names in this report does not constitute an official endorsement or approval of the use of such commercial hardware or software. The report may not be cited for purposes of advertisement.")

    # next page
    document.add_page_break()
    document.add_picture('images/Capture4.PNG', width=Inches(4))
    document.add_heading("Combat Capabilities Development Command (CCDC) Data &amp; Analysis Center (DAC) Enter "
                         "System Name Enter Event Type (e.g., CVPA, CVI, VoF, etc) Report", level=1)

    # get event team, get the analysts ann their first and last name here
    p = document.add_paragraph()
    p.add_run("by ...").bold = True

    document.add_paragraph("CCDC Data & Analysis Center")
    document.add_paragraph("Author(s)")
    document.add_paragraph("Affiliation")

    document.add_paragraph("classified by:" + event.getClassifiedBy())
    document.add_paragraph("Derived from:" + event.getSecurityClassificationTitleGuide())
    document.add_paragraph("Declassify on:" + event.getDeclassificationDate())

    # next page
    document.add_page_break()

    table = document.add_table(rows=31, cols=15)
    table.style = 'TableGrid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'REPORT DOCUMENTATION PAGE'

    a = table.cell(0, 0)
    b = table.cell(0, 9)
    a.merge(b)
    a = table.cell(0, 10)
    b = table.cell(0, 14)
    a.merge(b)

    hdr_cells = table.rows[1].cells
    hdr_cells[0].text = ''
    a = table.cell(1, 0)
    b = table.cell(1, 14)
    a.merge(b)

    hdr_cells = table.rows[2].cells
    hdr_cells[0].text = '1. REPORT DATE \n' + today.strftime("%B %d, %Y")
    a = table.cell(2, 0)
    b = table.cell(3, 4)
    a.merge(b)

    hdr_cells[5].text = '2. REPORT TYPE \n' + event.getType()
    a = table.cell(2, 5)
    b = table.cell(3, 10)
    a.merge(b)

    hdr_cells[11].text = '3. DATES COVERED'
    a = table.cell(2, 11)
    b = table.cell(3, 14)
    a.merge(b)

    hdr_cells = table.rows[4].cells
    hdr_cells[
        0].text = '4. TITLE AND SUBTITLE \n' + "Combat Capabilities Development Command (CCDC) Data &amp; Analysis Center (DAC) " + event.getType() + " Report."
    a = table.cell(4, 0)
    b = table.cell(6, 10)
    a.merge(b)
    hdr_cells[11].text = '5a. CONTRACT NUMBER'
    a = table.cell(4, 11)
    b = table.cell(4, 14)
    a.merge(b)

    hdr_cells = table.rows[5].cells
    hdr_cells[11].text = '5b. CONTRACT NUMBER'
    a = table.cell(5, 11)
    b = table.cell(5, 14)
    a.merge(b)

    hdr_cells = table.rows[6].cells
    hdr_cells[11].text = '5c. PROGRAM ELEMENT NUMBER'
    a = table.cell(6, 11)
    b = table.cell(6, 14)
    a.merge(b)

    hdr_cells = table.rows[7].cells
    hdr_cells[11].text = '5d. PROJECT NUMBER'
    a = table.cell(7, 11)
    b = table.cell(7, 14)
    a.merge(b)
    hdr_cells[0].text = '6. AUTHOR(S) \n' + my_string
    a = table.cell(7, 0)
    b = table.cell(10, 10)
    a.merge(b)

    hdr_cells = table.rows[8].cells
    hdr_cells[11].text = '5e. TASK NUMBER'
    a = table.cell(8, 11)
    b = table.cell(8, 14)
    a.merge(b)

    hdr_cells = table.rows[9].cells
    hdr_cells[11].text = '5f. WORK UNIT NUMBER'
    a = table.cell(9, 11)
    b = table.cell(9, 14)
    a.merge(b)

    hdr_cells = table.rows[10].cells
    hdr_cells[11].text = '8. PERFORMING ORGANIZATION REPORT NUMBER'
    a = table.cell(10, 11)
    b = table.cell(13, 14)
    a.merge(b)

    hdr_cells = table.rows[11].cells
    hdr_cells[
        0].text = '7. PERFORMING ORGANIZATION NAME(S) AND ADDRESS(ES) \n' + "Director U.S. Army CCDC Data & Analysis Center 6896 Mauchly Street Aberdeen Proving Ground, MD"
    a = table.cell(11, 0)
    b = table.cell(13, 10)
    a.merge(b)

    hdr_cells = table.rows[14].cells
    hdr_cells[0].text = '9. SPONSORING / MONITORING AGENCY NAME(S) AND ADDRESS(ES)'
    a = table.cell(14, 0)
    b = table.cell(17, 10)
    a.merge(b)
    hdr_cells[11].text = '10. SPONSOR/MONITOR’S ACRONYM(S)'
    a = table.cell(14, 11)
    b = table.cell(15, 14)
    a.merge(b)

    hdr_cells = table.rows[16].cells
    hdr_cells[11].text = '11. SPONSOR/MONITOR’S REPORT NUMBER(S)'
    a = table.cell(16, 11)
    b = table.cell(17, 14)
    a.merge(b)

    hdr_cells = table.rows[18].cells
    hdr_cells[0].text = '12. DISTRIBUTION / AVAILABILITY STATEMENT'
    a = table.cell(18, 0)
    b = table.cell(20, 14)
    a.merge(b)

    hdr_cells = table.rows[21].cells
    hdr_cells[0].text = '13. SUPPLEMENTARY NOTES'
    a = table.cell(21, 0)
    b = table.cell(23, 14)
    a.merge(b)

    hdr_cells = table.rows[23].cells
    hdr_cells[0].text = '14. ABSTRACT'
    a = table.cell(23, 0)
    b = table.cell(24, 14)
    a.merge(b)

    hdr_cells = table.rows[25].cells
    hdr_cells[0].text = '15. SUBJECT TERMS \n' + event.getType()
    a = table.cell(25, 0)
    b = table.cell(26, 14)
    a.merge(b)

    hdr_cells = table.rows[27].cells
    hdr_cells[0].text = '16. SECURITY CLASSIFICATION OF:'
    a = table.cell(27, 0)
    b = table.cell(28, 8)
    a.merge(b)

    hdr_cells[9].text = '17. LIMITATION OF ABSTRACT \n' + event.getEventClassification()
    a = table.cell(27, 9)
    b = table.cell(30, 10)
    a.merge(b)
    hdr_cells[11].text = '18. NUMBER OF PAGES \n' + "23"
    a = table.cell(27, 11)
    b = table.cell(30, 12)
    a.merge(b)
    hdr_cells[13].text = '19a. NAME OF RESPONSIBLE PERSON \n Person'
    a = table.cell(27, 13)
    b = table.cell(28, 14)
    a.merge(b)

    hdr_cells = table.rows[29].cells
    hdr_cells[0].text = 'a. REPORT \n' + event.getEventClassification()
    a = table.cell(29, 0)
    b = table.cell(30, 2)
    a.merge(b)
    hdr_cells[13].text = '19b. TELEPHONE NUMBER \n Phone'
    a = table.cell(29, 13)
    b = table.cell(30, 14)
    a.merge(b)
    hdr_cells[3].text = 'b. ABSTRACT \n' + event.getEventClassification()
    a = table.cell(29, 3)
    b = table.cell(30, 5)
    a.merge(b)
    hdr_cells[6].text = 'c. THIS PAGE \n' + event.getEventClassification()
    a = table.cell(29, 6)
    b = table.cell(30, 8)
    a.merge(b)

    # next page
    document.add_page_break()
    document.add_heading("Table of Contents")
    table = document.add_table(rows=15, cols=15)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'List of Figures'
    hdr_cells[14].text = 'iv'
    a = table.cell(0, 0)
    b = table.cell(0, 10)
    a.merge(b)

    hdr_cells = table.rows[1].cells
    hdr_cells[0].text = 'List of Tables'
    hdr_cells[14].text = 'v'
    a = table.cell(1, 0)
    b = table.cell(1, 10)
    a.merge(b)

    hdr_cells = table.rows[2].cells
    hdr_cells[0].text = 'Acknowledgements'
    hdr_cells[14].text = 'vi'
    a = table.cell(2, 0)
    b = table.cell(2, 10)
    a.merge(b)

    hdr_cells = table.rows[3].cells
    hdr_cells[0].text = 'Executive summary'
    hdr_cells[14].text = 'vii'
    a = table.cell(3, 0)
    b = table.cell(3, 10)
    a.merge(b)

    hdr_cells = table.rows[4].cells
    hdr_cells[0].text = '1. (U) INTRODUCTION'
    hdr_cells[14].text = '8'
    a = table.cell(4, 0)
    b = table.cell(4, 12)
    a.merge(b)

    hdr_cells = table.rows[5].cells
    hdr_cells[1].text = '1.1 (U) System/Network Architecture'
    hdr_cells[14].text = '9'
    a = table.cell(5, 1)
    b = table.cell(5, 10)
    a.merge(b)
    hdr_cells = table.rows[6].cells
    hdr_cells[1].text = '1.2 (U) Test Setup and Network Postures'
    hdr_cells[14].text = '10'
    a = table.cell(6, 1)
    b = table.cell(6, 10)
    a.merge(b)
    hdr_cells = table.rows[7].cells
    hdr_cells[1].text = '1.3 (U) Limitations'
    hdr_cells[14].text = '10'
    a = table.cell(7, 1)
    b = table.cell(7, 10)
    a.merge(b)

    hdr_cells = table.rows[8].cells
    hdr_cells[0].text = '2. ENTER EVENT TYPE (E.G., CVPA, CVI, VOF, ETC) FINDING'
    hdr_cells[14].text = '11'
    a = table.cell(8, 0)
    b = table.cell(8, 12)
    a.merge(b)

    hdr_cells = table.rows[9].cells
    hdr_cells[1].text = '2.1 (U) Lack of Encryption'
    hdr_cells[14].text = '12'
    a = table.cell(9, 1)
    b = table.cell(9, 10)
    a.merge(b)
    hdr_cells = table.rows[10].cells
    hdr_cells[1].text = '2.2 (U) Missing Patches'
    hdr_cells[14].text = '14'
    a = table.cell(10, 1)
    b = table.cell(10, 10)
    a.merge(b)

    hdr_cells = table.rows[11].cells
    hdr_cells[0].text = '3. CONCLUSIONS AND RECOMMENDATIONS'
    hdr_cells[14].text = '16'
    a = table.cell(11, 0)
    b = table.cell(11, 12)
    a.merge(b)

    hdr_cells = table.rows[12].cells
    hdr_cells[0].text = 'Appendix A - List of Acronym'
    hdr_cells[14].text = 'A-1'
    a = table.cell(12, 0)
    b = table.cell(12, 12)
    a.merge(b)
    hdr_cells = table.rows[13].cells
    hdr_cells[0].text = '3. Appendix B - Distribution List'
    hdr_cells[14].text = 'B-1'
    a = table.cell(13, 0)
    b = table.cell(13, 12)
    a.merge(b)

    # next page
    document.add_page_break()
    document.add_heading("List of Figures")

    # next page
    document.add_page_break()
    document.add_heading("List of Tables")
    table = document.add_table(rows=15, cols=15)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Table 1. (S) List of Findings'
    hdr_cells[14].text = '11'
    a = table.cell(0, 0)
    b = table.cell(0, 10)
    a.merge(b)
    hdr_cells = table.rows[1].cells
    hdr_cells[0].text = 'Table 2 describes the Lack of Encryption vulnerability.'
    hdr_cells[14].text = '12'
    a = table.cell(1, 0)
    b = table.cell(1, 10)
    a.merge(b)
    hdr_cells = table.rows[2].cells
    hdr_cells[0].text = 'Table 3 describes the missing Patches vulnerability.'
    hdr_cells[14].text = '14'
    a = table.cell(2, 0)
    b = table.cell(2, 10)
    a.merge(b)

    # next page
    document.add_page_break()
    document.add_heading("Acknowledgements", level=1)
    document.add_paragraph("The U.S Army Combat Capabilities De")
    # next page
    document.add_page_break()
    document.add_heading("Executive Summary", level=1)
    # next page
    document.add_page_break()
    document.add_heading("(U) INTRODUCTION", level=1)
    # next page
    document.add_page_break()
    document.add_heading("(U) System/Network Architecture", level=1)

    # next page
    document.add_page_break()
    document.add_paragraph("1.2 Test Setup and Network Postures")
    document.add_paragraph("Limitations")
    document.add_heading("(U) ENTER EVENT TYPE (E.G., CVPA, CVI, VOF, ETC) FINDINGS", level=1)
    document.add_paragraph("(U) Table 1 lists vulnerabilities identified and validated by CCDC DAC during the Enter "
                           "Event Type (e.g., CVPA, CVI, VoF, etc) with their associated technical risk.")
    document.add_heading("Table 1. (S) List of Findings", level=1)

    table = document.add_table(rows=1, cols=9)
    table.style = 'TableGrid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'ID'
    hdr_cells[1].text = 'DESCRIPTION'
    hdr_cells[6].text = 'LIKELIHOOD'
    hdr_cells[7].text = 'IMPACT'
    hdr_cells[8].text = 'RISK'

    row = 0
    a = table.cell(row, 1)
    b = table.cell(row, 5)
    a.merge(b)
    for finding in findingsList:
        table.add_row()
        row += 1
        hdr_cells = table.rows[row].cells
        hdr_cells[0].text = str(finding.getid())
        hdr_cells[1].text = finding.getDescription()
        hdr_cells[6].text = finding.getLikelihood()
        hdr_cells[7].text = str(finding.getImpactLevel())
        hdr_cells[8].text = finding.getRisk()
        a = table.cell(row, 1)
        b = table.cell(row, 5)
        a.merge(b)

    # next page
    document.add_page_break()
    document.add_heading("(U) Lack of Encryption")
    document.add_paragraph("Table 2 describes the Lack of Encryption vulnerability")
    document.add_heading("Table 2. Lack of Encryption")

    lackEnc = None
    for finding in findingsList:
        if finding.getType().value == 'Encryption':
            lackEnc = finding
        else:
            lackEnc = finding

    missPatch = None
    for finding in findingsList:
        if finding.getType().value == 'Missing Patches':
            missPatch = finding
        else:
            missPatch = finding

    # create a table for the finding with the encryption type
    table = document.add_table(rows=16, cols=12)
    table.style = 'TableGrid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'ID'
    hdr_cells[1].text = str(lackEnc.getid())
    a = table.cell(0, 1)
    b = table.cell(0, 2)
    a.merge(b)
    hdr_cells[3].text = ''

    hdr_cells[4].text = 'IMPACT SCORE'
    hdr_cells[5].text = str(lackEnc.getImpactScore())

    hdr_cells[6].text = 'STATUS'
    hdr_cells[7].text = lackEnc.getStatus()
    a = table.cell(0, 7)
    b = table.cell(0, 8)
    a.merge(b)
    a = table.cell(1, 7)
    b = table.cell(1, 8)
    a.merge(b)
    hdr_cells[9].text = 'POSTURE'
    a = table.cell(0, 9)
    b = table.cell(0, 11)
    a.merge(b)
    a = table.cell(1, 9)
    b = table.cell(1, 11)
    a.merge(b)
    #
    #
    hdr_cells = table.rows[1].cells
    hdr_cells[0].text = 'HOST NAMES'
    hdr_cells[2].text = 'IP:PORT'
    hdr_cells[4].text = 'CAT'
    hdr_cells[5].text = str(lackEnc.getSeverityCategoryCode().value)
    hdr_cells[6].text = 'LIKELIHOOD'
    hdr_cells[7].text = lackEnc.getLikelihood()
    hdr_cells[9].text = lackEnc.getPosture()

    a = table.cell(1, 0)
    b = table.cell(1, 1)
    a.merge(b)

    a = table.cell(1, 2)
    b = table.cell(1, 3)
    a.merge(b)
    a = table.cell(2, 2)
    b = table.cell(4, 3)
    a.merge(b)
    a = table.cell(2, 0)
    b = table.cell(4, 1)
    a.merge(b)
    a = table.cell(1, 9)
    b = table.cell(2, 11)
    a.merge(b)

    hdr_cells = table.rows[2].cells
    hdr_cells[0].text = lackEnc.getHostName()
    hdr_cells[2].text = lackEnc.getIpPort()
    hdr_cells[4].text = 'CAT SCORE'
    hdr_cells[5].text = str(lackEnc.getSeverityCategoryScore())
    hdr_cells[6].text = 'IMPACT'
    hdr_cells[7].text = str(lackEnc.getImpactLevel().value)
    a = table.cell(2, 7)
    b = table.cell(2, 8)
    a.merge(b)

    hdr_cells = table.rows[3].cells
    hdr_cells[4].text = 'Vs Score'
    hdr_cells[5].text = str(lackEnc.getSeverityCategoryScore())
    hdr_cells[6].text = 'Risk'
    hdr_cells[7].text = str(lackEnc.getImpactLevel().value)
    hdr_cells[9].text = 'C'
    hdr_cells[10].text = 'I'
    hdr_cells[11].text = 'A'

    a = table.cell(3, 7)
    b = table.cell(3, 8)
    a.merge(b)

    hdr_cells = table.rows[4].cells
    hdr_cells[4].text = 'Vs'
    hdr_cells[5].text = str(lackEnc.getSeverityCategoryCode().value)
    hdr_cells[6].text = 'CM'
    hdr_cells[7].text = str(lackEnc.getCountermeasureEffectivenessRating().value)
    hdr_cells[9].text = lackEnc.getConfidentiality()
    hdr_cells[10].text = lackEnc.getIntegrity()
    hdr_cells[11].text = lackEnc.getAvailability()

    a = table.cell(4, 7)
    b = table.cell(4, 8)
    a.merge(b)

    hdr_cells = table.rows[5].cells
    hdr_cells[0].text = 'TYPE'
    hdr_cells[4].text = 'Encryption'
    a = table.cell(5, 4)
    b = table.cell(5, 11)
    a.merge(b)
    a = table.cell(5, 0)
    b = table.cell(5, 3)
    a.merge(b)

    hdr_cells = table.rows[6].cells
    hdr_cells[0].text = ''
    hdr_cells[1].text = lackEnc.getDescription()
    a = table.cell(6, 1)
    b = table.cell(6, 11)
    a.merge(b)

    hdr_cells = table.rows[7].cells
    hdr_cells[0].text = 'DESCRIPTION'
    hdr_cells[1].text = lackEnc.getLongDescription()
    a = table.cell(7, 1)
    b = table.cell(9, 11)
    a.merge(b)
    a = table.cell(7, 0)
    b = table.cell(9, 0)
    a.merge(b)

    hdr_cells = table.rows[10].cells
    hdr_cells[0].text = ''
    hdr_cells[1].text = lackEnc.getMitigationBriefDescription()
    a = table.cell(7, 1)
    b = table.cell(7, 11)
    a.merge(b)
    a = table.cell(10, 1)
    b = table.cell(10, 11)
    a.merge(b)

    hdr_cells = table.rows[11].cells
    hdr_cells[0].text = 'MITIGATION'
    hdr_cells[1].text = lackEnc.getMitigationLongDescription()
    a = table.cell(11, 1)
    b = table.cell(13, 11)
    a.merge(b)
    a = table.cell(11, 0)
    b = table.cell(13, 0)
    a.merge(b)

    hdr_cells = table.rows[14].cells
    hdr_cells[0].text = 'REFERENCE'
    hdr_cells[4].text = 'FIGURE X'
    a = table.cell(14, 0)
    b = table.cell(14, 3)
    a.merge(b)
    a = table.cell(14, 4)
    b = table.cell(14, 11)
    a.merge(b)

    hdr_cells = table.rows[15].cells
    hdr_cells[0].text = 'C-CONFIDENTIALITY  I-INTEGRITY  A-AVAILABILITY  CM-COUNTERMEASURE'
    a = table.cell(15, 0)
    b = table.cell(15, 11)
    a.merge(b)

    # next page
    document.add_page_break()
    document.add_heading("2.2 (U) Missing Patches")
    document.add_paragraph("Table 3 describes the missing patches vulnerability")
    document.add_heading("Table 3. Missing Patches")

    # create a table for the finding with the missing patches type
    table = document.add_table(rows=16, cols=12)
    table.style = 'TableGrid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'ID'
    hdr_cells[1].text = str(missPatch.getid())
    a = table.cell(0, 1)
    b = table.cell(0, 2)
    a.merge(b)
    hdr_cells[3].text = ''

    hdr_cells[4].text = 'IMPACT SCORE'
    hdr_cells[5].text = str(missPatch.getImpactScore())

    hdr_cells[6].text = 'STATUS'
    hdr_cells[7].text = missPatch.getStatus()
    a = table.cell(0, 7)
    b = table.cell(0, 8)
    a.merge(b)
    a = table.cell(1, 7)
    b = table.cell(1, 8)
    a.merge(b)
    hdr_cells[9].text = 'POSTURE'
    a = table.cell(0, 9)
    b = table.cell(0, 11)
    a.merge(b)
    a = table.cell(1, 9)
    b = table.cell(1, 11)
    a.merge(b)
    #
    #
    hdr_cells = table.rows[1].cells
    hdr_cells[0].text = 'HOST NAMES'
    hdr_cells[2].text = 'IP:PORT'
    hdr_cells[4].text = 'CAT'
    hdr_cells[5].text = str(missPatch.getSeverityCategoryCode().value)
    hdr_cells[6].text = 'LIKELIHOOD'
    hdr_cells[7].text = missPatch.getLikelihood()
    hdr_cells[9].text = missPatch.getPosture()

    a = table.cell(1, 0)
    b = table.cell(1, 1)
    a.merge(b)

    a = table.cell(1, 2)
    b = table.cell(1, 3)
    a.merge(b)
    a = table.cell(2, 2)
    b = table.cell(4, 3)
    a.merge(b)
    a = table.cell(2, 0)
    b = table.cell(4, 1)
    a.merge(b)
    a = table.cell(1, 9)
    b = table.cell(2, 11)
    a.merge(b)

    hdr_cells = table.rows[2].cells
    hdr_cells[0].text = missPatch.getHostName()
    hdr_cells[2].text = missPatch.getIpPort()
    hdr_cells[4].text = 'CAT SCORE'
    hdr_cells[5].text = str(missPatch.getSeverityCategoryScore())
    hdr_cells[6].text = 'IMPACT'
    hdr_cells[7].text = str(missPatch.getImpactLevel().value)
    a = table.cell(2, 7)
    b = table.cell(2, 8)
    a.merge(b)

    hdr_cells = table.rows[3].cells
    hdr_cells[4].text = 'Vs Score'
    hdr_cells[5].text = str(missPatch.getSeverityCategoryScore())
    hdr_cells[6].text = 'Risk'
    hdr_cells[7].text = str(missPatch.getImpactLevel().value)
    hdr_cells[9].text = 'C'
    hdr_cells[10].text = 'I'
    hdr_cells[11].text = 'A'

    a = table.cell(3, 7)
    b = table.cell(3, 8)
    a.merge(b)

    hdr_cells = table.rows[4].cells
    hdr_cells[4].text = 'Vs'
    hdr_cells[5].text = str(missPatch.getSeverityCategoryCode().value)
    hdr_cells[6].text = 'CM'
    hdr_cells[7].text = str(missPatch.getCountermeasureEffectivenessRating().value)
    hdr_cells[9].text = missPatch.getConfidentiality()
    hdr_cells[10].text = missPatch.getIntegrity()
    hdr_cells[11].text = missPatch.getAvailability()

    a = table.cell(4, 7)
    b = table.cell(4, 8)
    a.merge(b)

    hdr_cells = table.rows[5].cells
    hdr_cells[0].text = 'TYPE'
    hdr_cells[4].text = 'Missing Patches'
    a = table.cell(5, 4)
    b = table.cell(5, 11)
    a.merge(b)
    a = table.cell(5, 0)
    b = table.cell(5, 3)
    a.merge(b)

    hdr_cells = table.rows[6].cells
    hdr_cells[0].text = ''
    hdr_cells[1].text = missPatch.getDescription()
    a = table.cell(6, 1)
    b = table.cell(6, 11)
    a.merge(b)

    hdr_cells = table.rows[7].cells
    hdr_cells[0].text = 'DESCRIPTION'
    hdr_cells[1].text = missPatch.getLongDescription()
    a = table.cell(7, 1)
    b = table.cell(9, 11)
    a.merge(b)
    a = table.cell(7, 0)
    b = table.cell(9, 0)
    a.merge(b)

    hdr_cells = table.rows[10].cells
    hdr_cells[0].text = ''
    hdr_cells[1].text = missPatch.getMitigationBriefDescription()
    a = table.cell(7, 1)
    b = table.cell(7, 11)
    a.merge(b)
    a = table.cell(10, 1)
    b = table.cell(10, 11)
    a.merge(b)

    hdr_cells = table.rows[11].cells
    hdr_cells[0].text = 'MITIGATION'
    hdr_cells[1].text = missPatch.getMitigationLongDescription()
    a = table.cell(11, 1)
    b = table.cell(13, 11)
    a.merge(b)
    a = table.cell(11, 0)
    b = table.cell(13, 0)
    a.merge(b)

    hdr_cells = table.rows[14].cells
    hdr_cells[0].text = 'REFERENCE'
    hdr_cells[4].text = 'FIGURE X'
    a = table.cell(14, 0)
    b = table.cell(14, 3)
    a.merge(b)
    a = table.cell(14, 4)
    b = table.cell(14, 11)
    a.merge(b)

    hdr_cells = table.rows[15].cells
    hdr_cells[0].text = 'C-CONFIDENTIALITY  I-INTEGRITY  A-AVAILABILITY  CM-COUNTERMEASURE'
    a = table.cell(15, 0)
    b = table.cell(15, 11)
    a.merge(b)

    # next page
    document.add_page_break()
    document.add_heading("3. CONCLUSIONS AND RECOMMENDATIONS")
    document.add_paragraph("Insert conclusions write up here")

    # next page
    document.add_page_break()
    document.add_heading("Appendix A - List of Acronym")

    # next page
    document.add_page_break()
    document.add_paragraph("CCDC           Combat Capabilities Development Command")
    document.add_paragraph("DAC            Data & Analysis Center")

    # next page
    document.add_page_break()
    document.add_heading("Appendix B - Distribution List")

    # next page
    document.add_page_break()
    document.add_heading("Organization")

    document.save(os.path.expanduser('~/Downloads/FinalTechnicalReport.docx'))


def createRiskMatrixReport(findingsList, event):
    # Create a new Excel file and add a RiskMatrixSheet, etc.
    workbook = xlsxwriter.Workbook(os.path.expanduser('~/Downloads/RiskMatrixReport.xlsx'))
    # sheets
    changeLogs = workbook.add_worksheet('CHANGE LOGS')
    RiskMatrixSheet = workbook.add_worksheet('RISK_MATRIX')
    histogram = workbook.add_worksheet('Histogram')
    overAllScoreCard = workbook.add_worksheet('OVERALL SCORE CARD')
    writeUpCards = workbook.add_worksheet('Write_Up_Cards')
    impact = workbook.add_worksheet('IMPACT')
    disaCatScores = workbook.add_worksheet('Disa Cat Scores')
    objectiveClassification = workbook.add_worksheet('OBJECTIVE CLASSIFICATION')
    countermeasure = workbook.add_worksheet('COUNTEMEASURE')
    tables = workbook.add_worksheet('TABLES')
    data = workbook.add_worksheet('DATA')
    likelihoodAndRisk = workbook.add_worksheet('LIKELIHOOD AND RISK')
    relevanceOfThreat = workbook.add_worksheet('RELEVANCE OF THREAT')

    # title format cell.
    tittle_format = workbook.add_format()
    tittle_format.set_font_size(40)
    tittle_format.set_bold(True)

    # title in tables format
    cell_format = workbook.add_format()
    cell_format.set_bg_color('#ADD8E6')
    cell_format.set_border(1)
    cell_format.set_border_color('#FFFFFF')
    cell_format.set_bold(True)
    cell_format.set_text_wrap(True)
    cell_format.set_align(alignment='left')

    # format for table content
    table_format = workbook.add_format()
    table_format.set_text_wrap(True)
    table_format.set_align(alignment='left')

    # format for table content
    cards_format = workbook.add_format()
    cards_format.set_text_wrap(True)
    cards_format.set_align('center')
    cards_format.set_align('vcenter')
    # title in tables format
    cards_title_format = workbook.add_format()
    cards_title_format.set_bg_color('#ADD8E6')
    cards_title_format.set_border(1)
    cards_title_format.set_border_color('#FFFFFF')
    cards_title_format.set_bold(True)
    cards_title_format.set_align('center')
    cards_title_format.set_align('vcenter')

    # title format cell.
    board_format = workbook.add_format()
    board_format.set_font_size(40)
    board_format.set_bold(True)
    board_format.set_bg_color('#ADD8E6')
    board_format.set_border(1)
    board_format.set_border_color('#FFFFFF')
    cards_title_format.set_align('center')
    cards_title_format.set_align('vcenter')

    # change logs sheet
    row = 1
    changeLogs.merge_range(row, 1, row, 2, "DATE", cell_format)
    changeLogs.merge_range(row, 3, row, 4, "VERSION", cell_format)
    changeLogs.merge_range(row, 5, row, 9, "CHANGES", cell_format)
    row += 1

    for logEntry in logsHandler.getAllLogs():
        changeLogs.merge_range(row, 1, row, 2, logEntry.getTime(), cards_format)
        changeLogs.merge_range(row, 3, row, 4, str(logEntry.getId()), cards_format)
        changeLogs.merge_range(row, 5, row, 9, logEntry.getAction(), cards_format)
        row += 1

    # Risk matrix sheet
    tagRow = 3
    # Write to file excel
    RiskMatrixSheet.write(1, 2, event.getName(), tittle_format)

    RiskMatrixSheet.write(tagRow, 1, "ID", cell_format)
    RiskMatrixSheet.set_column(1, 1, float(30))

    RiskMatrixSheet.write(tagRow, 2, "IP:PORT", cell_format)
    RiskMatrixSheet.write(tagRow, 3, "DESCRIPTION", cell_format)
    RiskMatrixSheet.set_column(3, 3, float(30))

    RiskMatrixSheet.write(tagRow, 4, "STATUS", cell_format)
    RiskMatrixSheet.write(tagRow, 5, "POSTURE", cell_format)
    RiskMatrixSheet.write(tagRow, 6, "C", cell_format)
    RiskMatrixSheet.write(tagRow, 7, "I", cell_format)
    RiskMatrixSheet.write(tagRow, 8, "A", cell_format)
    RiskMatrixSheet.write(tagRow, 9, "C", cell_format)
    RiskMatrixSheet.write(tagRow, 10, "I", cell_format)
    RiskMatrixSheet.write(tagRow, 11, "A", cell_format)
    RiskMatrixSheet.write(tagRow, 12, "IMP-SCORE", cell_format)
    RiskMatrixSheet.write(tagRow, 13, "CAT", cell_format)
    RiskMatrixSheet.write(tagRow, 14, "CAT SCORE", cell_format)
    RiskMatrixSheet.write(tagRow, 15, "CM", cell_format)
    RiskMatrixSheet.write(tagRow, 16, "Vs (n)", cell_format)
    RiskMatrixSheet.write(tagRow, 17, "Vs (Q)", cell_format)
    RiskMatrixSheet.write(tagRow, 18, "RELEVANCE OF THREAT", cell_format)
    RiskMatrixSheet.write(tagRow, 19, "LIKELIHOOD", cell_format)
    RiskMatrixSheet.write(tagRow, 20, "IMPACT", cell_format)
    # RiskMatrixSheet.writ(tagRow, 0, "IMPACT RATIONALE(HIDE IT)", bold)
    RiskMatrixSheet.write(tagRow, 21, "RISK", cell_format)

    #
    row = 4
    col = 1

    # Widen the first column to make the text clearer.
    RiskMatrixSheet.set_column('A:A', 20)

    for f in findingsList:
        RiskMatrixSheet.write(row, col, str(f.getid()), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getIpPort(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getDescription(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getStatus(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getPosture(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, "?", table_format)
        col += 1
        RiskMatrixSheet.write(row, col, "?", table_format)
        col += 1
        RiskMatrixSheet.write(row, col, "?", table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getConfidentiality(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getIntegrity(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getAvailability(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getImpactScore(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getSeverityCategoryCode(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getSeverityCategoryScore(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getCountermeasureEffectivenessRating(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getVulnerabilitySeverity(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getQualitativeVulnerabilitySeverity(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getRelevance(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getLikelihood(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getImpactLevel(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getRisk(), table_format)
        col += 1

        # for the small table
        RiskMatrixSheet.write(row, col, f.getIntegrityFindingImpactOnSystem(), table_format)
        col += 1
        RiskMatrixSheet.write(row, col, f.getAvailabilityFindingImpactOnSystem(), table_format)
        col += 1

        # update row and column
        row += 1
        col = 1

    # overall board sheet
    overAllScoreCard.merge_range(1, 2, 5, 20, "CYBER SCORE CARD", cards_title_format)
    overAllScoreCard.merge_range(7, 2, 8, 4, "FINDING", cards_title_format)
    overAllScoreCard.merge_range(7, 5, 8, 5, "RISK LEVEL", cards_title_format)

    overAllScoreCard.merge_range(7, 7, 7, 10, "CONFIDENTIALITY", cards_title_format)
    overAllScoreCard.merge_range(8, 7, 8, 8, "AFFECTED", cards_title_format)
    overAllScoreCard.merge_range(8, 9, 8, 10, "CV", cards_title_format)

    overAllScoreCard.merge_range(7, 12, 7, 15, "INTEGRITY", cards_title_format)
    overAllScoreCard.merge_range(8, 12, 8, 13, "AFFECTED", cards_title_format)
    overAllScoreCard.merge_range(8, 14, 8, 15, "IV", cards_title_format)

    overAllScoreCard.merge_range(7, 17, 7, 20, "AVAILABILITY", cards_title_format)
    overAllScoreCard.merge_range(8, 17, 8, 18, "AFFECTED", cards_title_format)
    overAllScoreCard.merge_range(8, 19, 8, 20, "AV", cards_title_format)
    row = 9
    for f in findingsList:
        if f.getConfidentiality() == 'Low':
            affected = 'NO'
        else:
            affected = 'YES'
        if f.getIntegrity() == 'Low':
            affected = 'NO'
        else:
            affected = 'YES'
        if f.getAvailability() == 'Low':
            affected1 = 'NO'
        else:
            affected1 = 'YES'
        overAllScoreCard.merge_range(row, 2, row, 4, f.getDescription())
        overAllScoreCard.write(row, 5, f.getRisk())

        overAllScoreCard.merge_range(row, 7, row, 8, affected)
        overAllScoreCard.merge_range(row, 9, row, 10, f.getConfidentiality())

        overAllScoreCard.merge_range(row, 12, row, 13, affected1)
        overAllScoreCard.merge_range(row, 14, row, 15, f.getIntegrity())

        overAllScoreCard.merge_range(row, 17, row, 18, affected)
        overAllScoreCard.merge_range(row, 19, row, 20, f.getAvailability())

        row += 1

    # write up cards sheet
    one = 1
    two = 2
    three = 3
    four = 4
    five = 5
    six = 6
    seven = 7
    eight = 8
    nine = 9
    ten = 10
    eleven = 11

    for f in findingsList:
        writeUpCards.write(3, one, "ID", cards_title_format)
        writeUpCards.write(3, two, str(f.getid()))
        writeUpCards.set_column(two, two, float(25))

        writeUpCards.merge_range(4, one, 4, two, "HOST NAMES", cards_title_format)
        writeUpCards.merge_range(5, one, 8, two, f.getHostName(), cards_format)

        writeUpCards.merge_range(4, three, 4, four, "IP:PORT", cards_title_format)
        writeUpCards.merge_range(5, three, 8, four, f.getIpPort(), cards_format)

        writeUpCards.set_column(five, five, float(20))

        writeUpCards.write(3, five, "IMPACT SCORE", cards_title_format)
        writeUpCards.write(3, six, f.getImpactScore(), cards_format)

        writeUpCards.write(4, five, "CAT", cards_title_format)
        writeUpCards.write(4, six, f.getSeverityCategoryCode(), cards_format)

        writeUpCards.write(5, five, "CAT SCORE", cards_title_format)
        writeUpCards.write(5, six, f.getSeverityCategoryScore(), cards_format)

        writeUpCards.write(6, five, "VS-SCORE", cards_title_format)
        writeUpCards.write(6, six, f.getQualitativeVulnerabilitySeverity(), cards_format)

        writeUpCards.write(7, five, "VS", cards_title_format)
        writeUpCards.write(7, six, f.getVulnerabilitySeverity(), cards_format)

        writeUpCards.write(8, five, "IMPACT RATIONALE", cards_title_format)
        writeUpCards.merge_range(8, six, 8, eleven, "impact rationale?", cards_format)

        # make rows 8,10,11 taller
        writeUpCards.set_row(8, 40)
        writeUpCards.set_row(10, 40)
        writeUpCards.set_row(11, 40)
        writeUpCards.set_column(seven, seven, float(12))

        writeUpCards.write(3, seven, "STATUS", cards_title_format)
        writeUpCards.write(3, eight, f.getStatus(), cards_format)
        writeUpCards.write(4, seven, "LIKELIHOOD", cards_title_format)
        writeUpCards.write(4, eight, f.getLikelihood(), cards_format)
        writeUpCards.write(5, seven, "IMPACT", cards_title_format)
        writeUpCards.write(5, eight, f.getImpactLevel(), cards_format)
        writeUpCards.write(6, seven, "RISK", cards_title_format)
        writeUpCards.write(6, eight, f.getRisk(), cards_format)
        writeUpCards.write(7, seven, "CM", cards_title_format)
        writeUpCards.write(7, eight, f.getCountermeasureEffectivenessRating(), cards_format)

        writeUpCards.merge_range(3, nine, 3, eleven, "POSTURE", cards_title_format)
        writeUpCards.merge_range(4, nine, 5, eleven, f.getPosture(), cards_format)

        writeUpCards.write(6, nine, "C", cards_title_format)
        writeUpCards.write(7, nine, f.getConfidentiality(), cards_format)
        writeUpCards.write(6, ten, "I", cards_title_format)
        writeUpCards.write(7, ten, f.getIntegrity(), cards_format)
        writeUpCards.write(6, eleven, "A", cards_title_format)
        writeUpCards.write(7, eleven, f.getAvailability(), cards_format)

        writeUpCards.merge_range(9, one, 9, two, "TYPE", cards_title_format)
        writeUpCards.merge_range(9, three, 9, eleven, f.getType(), cards_format)

        writeUpCards.merge_range(10, one, 10, two, "", cards_title_format)
        writeUpCards.merge_range(10, three, 10, eleven, f.getDescription(), cards_format)
        writeUpCards.merge_range(11, one, 11, two, "DESCRIPTION", cards_title_format)
        writeUpCards.merge_range(11, three, 11, eleven, f.getLongDescription(), cards_format)

        # update cols for the next finding
        one += 12
        two += 12
        three += 12
        four += 12
        five += 12
        six += 12
        seven += 12
        eight += 12
        nine += 12
        ten += 12
        eleven += 12

    # Impact sheet
    tagRow = 3
    # Write to file excel
    impact.merge_range(1, 1, 1, 8,
                       "(U) Impact is derived from the security objective (Confidentiality, Integrity, and Availability), the security objective rank (High, Medium, Low) and the number (n) of security objectives that are impacted. Semi-quantitate values are assigned to each vulnerability discovered using the conversion table shown in Table X3.",
                       table_format)
    impact.set_row(1, float(100))

    impact.merge_range(3, 1, 3, 8, "Table X3 (U) Impact conversion table", table_format)

    impact.merge_range(4, 1, 4, 3, "Number of Security Objectives Impacted (n)", table_format)
    impact.set_row(4, float(50))
    impact.set_row(5, float(50))
    impact.set_row(6, float(50))
    impact.set_row(7, float(50))
    impact.set_row(8, float(50))
    impact.set_row(9, float(50))
    impact.set_row(10, float(50))
    impact.set_row(11, float(50))
    impact.set_row(12, float(50))
    impact.set_row(13, float(50))

    impact.write(4, 4, "Objective Classification (High, Moderate, Low)", table_format)

    impact.write(4, 5, "Semi-Quantitative Value", table_format)
    impact.merge_range(4, 6, 4, 8, "Description", table_format)

    # fill table
    impact.merge_range(5, 1, 5, 3, 3, table_format)
    impact.merge_range(6, 1, 6, 3, 2, table_format)
    impact.merge_range(7, 1, 7, 3, 1, table_format)
    impact.merge_range(8, 1, 8, 3, 3, table_format)
    impact.merge_range(9, 1, 9, 3, 2, table_format)
    impact.merge_range(10, 1, 10, 3, 1, table_format)
    impact.merge_range(11, 1, 11, 3, 3, table_format)
    impact.merge_range(12, 1, 12, 3, 2, table_format)
    impact.merge_range(13, 1, 13, 3, 1, table_format)

    impact.write(5, 4, "HHH", table_format)
    impact.write(6, 4, "HHx", table_format)
    impact.write(7, 4, "Hxx", table_format)
    impact.write(8, 4, "MMM", table_format)
    impact.write(9, 4, "MMx", table_format)
    impact.write(10, 4, "Mxx", table_format)
    impact.write(11, 4, "LLL", table_format)
    impact.write(12, 4, "LLx", table_format)
    impact.write(13, 4, "Lxx", table_format)

    impact.write(5, 5, 10, table_format)
    impact.write(6, 5, 9, table_format)
    impact.write(7, 5, 8, table_format)
    impact.write(8, 5, 7, table_format)
    impact.write(9, 5, 6, table_format)
    impact.write(10, 5, 5, table_format)
    impact.write(11, 5, 4, table_format)
    impact.write(12, 5, 3, table_format)
    impact.write(13, 5, 2, table_format)

    impact.merge_range(5, 6, 5, 8, "Three security objectives are impacted and all three are ranked High", table_format)
    impact.merge_range(6, 6, 6, 8, "Two security objectives are impacted and both objectives are ranked High",
                       table_format)
    impact.merge_range(7, 6, 7, 8, "One security objective is impacted and the objective is ranked High", table_format)
    impact.merge_range(8, 6, 8, 8, "Three security objectives are impacted and all three are ranked Moderate",
                       table_format)
    impact.merge_range(9, 6, 9, 8, "Two security objectives are impacted and both objectives are ranked Moderate",
                       table_format)
    impact.merge_range(10, 6, 10, 8, "One security objective is impacted and the objective is ranked Moderate",
                       table_format)
    impact.merge_range(11, 6, 11, 8, "Three security objectives are impacted and all three are ranked Low",
                       table_format)
    impact.merge_range(12, 6, 12, 8, "Two security objectives are impacted and both objectives are ranked Low",
                       table_format)
    impact.merge_range(13, 6, 13, 8, "One security objective is impacted and the objective is ranked Low", table_format)

    # disa cat scores sheet
    disaCatScores.merge_range(1, 1, 1, 2, "Table X1 (U) DISA CAT Codes", table_format)
    disaCatScores.set_column(1, 1, float(10))
    disaCatScores.set_column(1, 2, float(30))

    disaCatScores.write(2, 1, "CAT Code", table_format)
    disaCatScores.write(2, 2, "DISA Severity Code Guideline", table_format)
    disaCatScores.write(3, 1, "CAT I", table_format)
    disaCatScores.write(3, 2,
                        "Any vulnerability, the exploitation of which will, directly and immediately result in the loss of Confidentiality, Integrity, or Availability",
                        table_format)
    disaCatScores.write(4, 1, "CAT II", table_format)
    disaCatScores.write(4, 2,
                        "Any vulnerability, the exploitation of which has a potential to result in loss of Confidentiality, Integrity or Availability",
                        table_format)
    disaCatScores.write(5, 1, "CAT III", table_format)
    disaCatScores.write(5, 2,
                        "Any vulnerability, the existence of which degrades a measure to protect against loss of Confidentiality, Integrity or Availability",
                        table_format)
    disaCatScores.merge_range(6, 1, 6, 2, "UNCLASSIFIED", table_format)

    disaCatScores.merge_range(8, 1, 8, 2,
                              " (U) The qualitative DISA CAT code for each vulnerability is then converted to a quantitative value using the conversion table shown in Table X2.",
                              table_format)

    disaCatScores.merge_range(10, 1, 10, 2, "Table X2 (U) DISA CAT Codes Conversion Table", table_format)
    disaCatScores.write(11, 1, "Qualitative Value", table_format)
    disaCatScores.write(11, 2, "Semi-Quantitative Value", table_format)
    disaCatScores.write(12, 1, "CAT I", table_format)
    disaCatScores.write(12, 2, 10, table_format)
    disaCatScores.write(13, 1, "CAT II", table_format)
    disaCatScores.write(13, 2, 7, table_format)
    disaCatScores.write(14, 1, "CAT III", table_format)
    disaCatScores.write(14, 2, 4, table_format)

    # objective classification sheet
    objectiveClassification.merge_range(1, 1, 1, 5,
                                        " (U) The Objective Classification was defined in the program’s RMF package as shown in Table X4:",
                                        table_format)
    objectiveClassification.set_column(1, 2, float(30))

    objectiveClassification.merge_range(3, 1, 3, 5, "Table X4 (U) Objective Classification", table_format)
    objectiveClassification.merge_range(4, 1, 4, 2, "Security Objective", table_format)
    objectiveClassification.write(4, 3, "Confidentiality (C)", table_format)
    objectiveClassification.write(4, 4, "Integrity (I)", table_format)
    objectiveClassification.write(4, 5, "Availability (A)", table_format)

    objectiveClassification.merge_range(5, 1, 5, 2, "Objective Rank", table_format)
    objectiveClassification.write(5, 3, "Moderate", table_format)
    objectiveClassification.write(5, 4, "High", table_format)
    objectiveClassification.write(5, 5, "Moderate", table_format)

    objectiveClassification.merge_range(6, 1, 6, 5, "U/FOUO", table_format)

    objectiveClassification.merge_range(8, 1, 8, 2,
                                        "(U//FOUO) From the Table X4, any vulnerability that impacts Integrity will be assigned a semi-quantitative value of 8 for Impact. Any vulnerability that affects both Confidentiality and Availability will be assigned a semi-quantitative value of 6 for Impact. Any vulnerability that affects either Confidentiality or Availability will be assigned a semi-quantitative value of 5 for Impact.",
                                        table_format)
    objectiveClassification.set_row(8, float(100))

    objectiveClassification.merge_range(10, 1, 10, 2, "UNCLASSIFIED", table_format)
    objectiveClassification.merge_range(11, 1, 11, 2,
                                        " (U) Vulnerability Severity (Vs) is calculated using the formula in Fig. X0. The resulting quantitative value is then converted to a qualitative value using Table X6.",
                                        table_format)
    objectiveClassification.set_row(11, float(60))

    objectiveClassification.merge_range(11, 1, 11, 2, "Table X6 (U) Assessment scale", table_format)
    objectiveClassification.write(12, 1, "Qualitative Values", table_format)
    objectiveClassification.write(13, 1, "Very High", table_format)
    objectiveClassification.write(14, 1, "High", table_format)
    objectiveClassification.write(15, 1, "Moderate", table_format)
    objectiveClassification.write(16, 1, "Low", table_format)
    objectiveClassification.write(17, 1, "Very Low", table_format)
    objectiveClassification.merge_range(18, 1, 17, 2, "UNCLASSIFIED", table_format)

    objectiveClassification.write(12, 1, "Semi-Quantitative Values", table_format)
    objectiveClassification.write(13, 2, "95 ≤ VS ≤ 100", table_format)
    objectiveClassification.write(14, 2, "80 ≤ VS < 95", table_format)
    objectiveClassification.write(15, 2, "20 ≤ VS < 80", table_format)
    objectiveClassification.write(16, 2, "5 ≤ VS < 20", table_format)
    objectiveClassification.write(17, 2, "0 ≤ VS < 5", table_format)

    objectiveClassification.merge_range(1, 8, 1, 13,
                                        "(U) As an example, the Vulnerability Severity of a vulnerability discovered during the assessment may be calculated as:",
                                        table_format)
    objectiveClassification.merge_range(2, 8, 2, 13, "(U) Sraw (DISA CAT I finding) = 10 ", table_format)
    objectiveClassification.merge_range(3, 8, 3, 13, "(U) Impact (vulnerability only impacts Integrity) = 8",
                                        table_format)
    objectiveClassification.merge_range(4, 8, 4, 13, "(U) Countermeasure (countermeasure partially in place) = 5",
                                        table_format)
    objectiveClassification.merge_range(5, 8, 5, 13, "(U) Vs = (10*8*5)/10 = 40)", table_format)
    objectiveClassification.merge_range(6, 8, 5, 13,
                                        "(U) Using Table X6, and the value of 40 for the Vs, the qualitative value of the Vulnerability Severity is assessed as Moderate. ",
                                        table_format)
    objectiveClassification.merge_range(7, 8, 7, 13, "(U) Relevance of Threat is defined in Table X7.", table_format)

    # Countermeasure sheet
    countermeasure.merge_range(1, 2, 1, 10,
                               "(U) Countermeasure is assessed on a scale form 0-10 and determines the effectiveness of a countermeasure against a discovered vulnerability. The countermeasure value is derived from Table X5.",
                               table_format)
    countermeasure.set_row(1, float(60))
    countermeasure.set_column(1, 2, float(20))

    countermeasure.merge_range(3, 2, 3, 10, "Table X5 (U) Countermeasure Table", table_format)
    countermeasure.write(4, 2, "Qualitative Value", table_format)
    countermeasure.merge_range(4, 3, 4, 4, "Semi-Quantitative Value", table_format)
    countermeasure.merge_range(4, 5, 4, 10, "Description", table_format)

    countermeasure.write(5, 1, "Sucks", table_format)
    countermeasure.write(9, 1, "Great", table_format)

    countermeasure.write(5, 2, "Very High", table_format)
    countermeasure.write(6, 2, "High", table_format)
    countermeasure.write(7, 2, "Moderate", table_format)
    countermeasure.write(8, 2, "Low", table_format)
    countermeasure.write(9, 2, "Very Low", table_format)

    countermeasure.merge_range(5, 3, 5, 4, 10, table_format)
    countermeasure.merge_range(6, 3, 6, 4, "7-9", table_format)
    countermeasure.merge_range(7, 3, 7, 4, "4-6", table_format)
    countermeasure.merge_range(8, 3, 8, 4, "1--3", table_format)
    countermeasure.merge_range(9, 3, 9, 4, 0, table_format)

    countermeasure.merge_range(5, 5, 5, 10, "Countermeasure not implemented ", table_format)
    countermeasure.merge_range(6, 5, 6, 10, "Countermeasure is implemented but MINIMALLY effective", table_format)
    countermeasure.merge_range(7, 5, 7, 10, "Countermeasure is implemented but MODERATELY effective", table_format)
    countermeasure.merge_range(8, 5, 8, 10, "Countermeasure is implemented HIGHLY effective but can be improved.",
                               table_format)
    countermeasure.merge_range(9, 5, 9, 10, "Countermeasure implemented is effective", table_format)

    countermeasure.merge_range(5, 11, 5, 15, "It's absent, there is nothing that stops you from exploiting",
                               table_format)
    countermeasure.merge_range(6, 11, 6, 15, "Requires minimal effort to bypass", table_format)
    countermeasure.merge_range(7, 11, 7, 15, "Requires medium level of effort to bypass", table_format)
    countermeasure.merge_range(8, 11, 8, 15, "Requires a high level of effort  to bypass.", table_format)
    countermeasure.merge_range(9, 11, 9, 15, "You could not bypass the control, you had to get whitecarded",
                               table_format)

    countermeasure.merge_range(11, 2, 11, 5, "The following values are recommended.", table_format)
    countermeasure.merge_range(12, 2, 12, 5, "PORT SECURITY", table_format)

    countermeasure.merge_range(13, 2, 13, 3, "ZERO PORT SECURITY", table_format)
    countermeasure.merge_range(14, 2, 14, 3, "ADMIN PORTS DISABLED", table_format)
    countermeasure.merge_range(15, 2, 15, 3, "port security-dynamic", table_format)
    countermeasure.merge_range(16, 2, 16, 3, "port security static", table_format)
    countermeasure.merge_range(17, 2, 17, 3, "802.1X  AUTHENTICATION", table_format)

    countermeasure.write(13, 5, 10, table_format)
    countermeasure.write(14, 5, 8, table_format)
    countermeasure.write(15, 5, 5, table_format)
    countermeasure.write(16, 5, 2, table_format)
    countermeasure.write(17, 5, "0 UNLESS YOU BYPASS", table_format)

    countermeasure.merge_range(19, 5, 19, 10,
                               "Only give credit to the controls assessed, for example, if you did not (stress/tested/assessed)  the lock on the racks, don't included.",
                               table_format)
    countermeasure.merge_range(20, 5, 20, 10,
                               "!!!!But you must include the context on the report, what level of access was granted in order for you to conduct the assessment!!!!.",
                               table_format)
    countermeasure.merge_range(21, 5, 21, 10, "If you rate it a Low, give rationale and how it can be improved",
                               table_format)
    countermeasure.set_row(19, float(60))
    countermeasure.set_row(20, float(60))
    countermeasure.set_row(21, float(60))

    # tables sheet
    tables.merge_range(1, 1, 1, 10, "Likelihood", cards_title_format)
    tables.set_column(1, 1, float(10))
    tables.set_column(1, 2, float(10))
    tables.set_column(1, 3, float(10))
    tables.set_column(1, 4, float(10))
    tables.set_column(1, 5, float(10))
    tables.set_column(1, 6, float(10))
    tables.set_column(1, 7, float(10))
    tables.merge_range(2, 1, 9, 1, "Relevance of Threat", table_format)
    tables.merge_range(2, 2, 2, 7, "Vulnerability Severity (Vs)", table_format)

    tables.write(3, 3, "VL", table_format)
    tables.write(3, 4, "L", table_format)
    tables.write(3, 5, "M", table_format)
    tables.write(3, 6, "H", table_format)
    tables.write(3, 7, "VH", table_format)
    tables.write(3, 8, "INFO", table_format)

    tables.write(4, 3, 2, table_format)
    tables.write(4, 4, 3, table_format)
    tables.write(4, 5, 4, table_format)
    tables.write(4, 6, 5, table_format)
    tables.write(4, 7, 6, table_format)

    tables.write(5, 2, "Confirmed", table_format)
    tables.write(6, 2, "Expected", table_format)
    tables.write(7, 2, "Anticipated", table_format)
    tables.write(8, 2, "Predicted", table_format)
    tables.write(9, 2, "Possible", table_format)

    tables.write(5, 3, "VL", table_format)
    tables.write(6, 3, "VL", table_format)
    tables.write(7, 3, "VL", table_format)
    tables.write(8, 3, "VL", table_format)
    tables.write(9, 3, "VL", table_format)

    tables.write(5, 4, "L", table_format)
    tables.write(6, 4, "L", table_format)
    tables.write(7, 4, "L", table_format)
    tables.write(8, 4, "L", table_format)
    tables.write(9, 4, "VL", table_format)

    tables.write(5, 5, "M", table_format)
    tables.write(6, 5, "M", table_format)
    tables.write(7, 5, "M", table_format)
    tables.write(8, 5, "L", table_format)
    tables.write(9, 5, "L", table_format)

    tables.write(5, 6, "H", table_format)
    tables.write(6, 6, "H", table_format)
    tables.write(7, 6, "M", table_format)
    tables.write(8, 6, "L", table_format)
    tables.write(9, 6, "L", table_format)

    tables.write(5, 7, "VH", table_format)
    tables.write(6, 7, "VH", table_format)
    tables.write(7, 7, "H", table_format)
    tables.write(8, 7, "M", table_format)
    tables.write(9, 7, "L", table_format)

    tables.write(5, 8, "INFO", table_format)
    tables.write(6, 8, "INFO", table_format)
    tables.write(7, 8, "INFO", table_format)
    tables.write(8, 8, "INFO", table_format)
    tables.write(9, 8, "INFO", table_format)

    # SECOND TABLE IN THIS SHEET
    tables.merge_range(12, 1, 12, 10, "Assessed Risk", cards_title_format)

    tables.merge_range(13, 1, 21, 1, "Relevance of Threat", table_format)
    tables.merge_range(13, 2, 13, 7, "ImpactTier3", table_format)

    tables.write(14, 3, "VL", table_format)
    tables.write(14, 4, "L", table_format)
    tables.write(14, 5, "M", table_format)
    tables.write(14, 6, "H", table_format)
    tables.write(14, 7, "VH", table_format)
    tables.write(14, 8, "INFO", table_format)

    tables.write(15, 3, 2, table_format)
    tables.write(15, 4, 3, table_format)
    tables.write(15, 5, 4, table_format)
    tables.write(15, 6, 5, table_format)
    tables.write(15, 7, 6, table_format)

    tables.write(16, 2, "VH", table_format)
    tables.write(17, 2, "H", table_format)
    tables.write(18, 2, "M", table_format)
    tables.write(19, 2, "L", table_format)
    tables.write(20, 2, "VL", table_format)
    tables.write(21, 2, "INFO", table_format)

    tables.write(16, 3, "VL", table_format)
    tables.write(17, 3, "VL", table_format)
    tables.write(18, 3, "VL", table_format)
    tables.write(19, 3, "VL", table_format)
    tables.write(20, 3, "VL", table_format)
    tables.write(21, 3, "INFO", table_format)

    tables.write(16, 4, "L", table_format)
    tables.write(17, 4, "L", table_format)
    tables.write(18, 4, "L", table_format)
    tables.write(19, 4, "L", table_format)
    tables.write(20, 4, "VL", table_format)
    tables.write(21, 4, "INFO", table_format)

    tables.write(16, 5, "M", table_format)
    tables.write(17, 5, "M", table_format)
    tables.write(18, 5, "M", table_format)
    tables.write(19, 5, "L", table_format)
    tables.write(20, 5, "L", table_format)
    tables.write(21, 5, "INFO", table_format)

    tables.write(16, 6, "H", table_format)
    tables.write(17, 6, "H", table_format)
    tables.write(18, 6, "M", table_format)
    tables.write(19, 6, "L", table_format)
    tables.write(20, 6, "L", table_format)
    tables.write(21, 6, "INFO", table_format)

    tables.write(16, 7, "VH", table_format)
    tables.write(17, 7, "VH", table_format)
    tables.write(18, 7, "H", table_format)
    tables.write(19, 7, "M", table_format)
    tables.write(20, 7, "L", table_format)
    tables.write(21, 7, "INFO", table_format)

    tables.write(16, 8, "INFO", table_format)
    tables.write(17, 8, "INFO", table_format)
    tables.write(18, 8, "INFO", table_format)
    tables.write(19, 8, "INFO", table_format)
    tables.write(20, 8, "INFO", table_format)

    # Data sheet
    data.write(2, 2, "LEVEL", cards_title_format)
    data.write(3, 2, "M", table_format)
    data.write(4, 2, "H", table_format)
    data.write(5, 2, "L", table_format)

    data.write(2, 4, "STATUS", cards_title_format)
    data.write(3, 4, "OPEN", table_format)
    data.write(4, 4, "CLOSED", table_format)

    data.write(2, 6, "POSTURE", cards_title_format)
    data.write(3, 6, "INSIDER-1", table_format)
    data.write(4, 6, "INSIDER-2", table_format)
    data.write(5, 6, "INSIDER-3", table_format)
    data.write(6, 6, "OUTSIDER-1", table_format)
    data.write(7, 6, "OUTSIDER-2", table_format)
    data.write(8, 6, "OUTSIDER-3", table_format)
    data.write(9, 6, "NEARSIDER-1", table_format)
    data.write(10, 6, "NEARSIDER-2", table_format)
    data.write(11, 6, "NEARSIDER-3", table_format)

    data.write(2, 8, "CAT", cards_title_format)
    data.write(3, 8, "I", table_format)
    data.write(4, 8, "II", table_format)
    data.write(5, 8, "III", table_format)

    data.write(2, 10, "RELEVANCE", cards_title_format)
    data.write(3, 10, "Confirmed", table_format)
    data.write(4, 10, "Expected", table_format)
    data.write(5, 10, "Anticipated", table_format)
    data.write(6, 10, "Predicted", table_format)
    data.write(7, 10, "Possible", table_format)

    data.write(2, 12, "LEVELS", cards_title_format)
    data.write(3, 12, "VH", table_format)
    data.write(4, 12, "H", table_format)
    data.write(5, 12, "M", table_format)
    data.write(6, 12, "L", table_format)
    data.write(7, 12, "VL", table_format)
    data.write(8, 12, "INFO", table_format)

    data.set_column(1, 2, float(30))

    data.write(9, 2, "FINDING CATERGORIES", cards_title_format)
    data.write(10, 2, "CREDENTIALS COMPLEXITY", table_format)
    data.write(11, 2, "MANUFACTURER DEFAULT CREDS", table_format)
    data.write(12, 2, "LACK OF AUTHENTICATION", table_format)
    data.write(13, 2, "PLAIN TEXT PROTOCOLS", table_format)
    data.write(14, 2, "PLAIN TEXT WEB-LOGIN", table_format)
    data.write(15, 2, "ENCRYPTION", table_format)
    data.write(16, 2, "AUTHENTICATION BYPASS", table_format)
    data.write(17, 2, "PORT SECURITY", table_format)
    data.write(18, 2, "ACCESS CONTROL", table_format)
    data.write(19, 2, "LEAST PRIVILEGE", table_format)
    data.write(20, 2, "PRIVILEGE ESCALATION", table_format)
    data.write(21, 2, "MISSING PATCHES", table_format)
    data.write(22, 2, "PHYSICAL SECURITY", table_format)

    data.merge_range(24, 2, 24, 4, "System Categorization", cards_title_format)
    data.write(25, 2, "Confidentiality", cards_title_format)
    data.write(25, 3, "Integrity", cards_title_format)
    data.write(25, 4, "Availability", cards_title_format)

    data.write(27, 2, "H", table_format)
    data.write(28, 2, "H", table_format)
    data.write(29, 2, "H", table_format)
    data.write(30, 2, "M", table_format)
    data.write(31, 2, "M", table_format)
    data.write(32, 2, "M", table_format)
    data.write(33, 2, "L", table_format)
    data.write(34, 2, "L", table_format)
    data.write(35, 2, "L", table_format)

    data.write(27, 3, "H", table_format)
    data.write(28, 3, "H", table_format)
    data.write(29, 3, "X", table_format)
    data.write(30, 3, "M", table_format)
    data.write(31, 3, "M", table_format)
    data.write(32, 3, "X", table_format)
    data.write(33, 3, "L", table_format)
    data.write(34, 3, "L", table_format)
    data.write(35, 3, "X", table_format)

    data.write(27, 4, "H", table_format)
    data.write(28, 4, "X", table_format)
    data.write(29, 4, "X", table_format)
    data.write(30, 4, "M", table_format)
    data.write(31, 4, "X", table_format)
    data.write(32, 4, "X", table_format)
    data.write(33, 4, "L", table_format)
    data.write(34, 4, "X", table_format)
    data.write(35, 4, "X", table_format)

    data.write(27, 5, 10, table_format)
    data.write(28, 5, 9, table_format)
    data.write(29, 5, 8, table_format)
    data.write(30, 5, 7, table_format)
    data.write(31, 5, 6, table_format)
    data.write(32, 5, 5, table_format)
    data.write(33, 5, 4, table_format)
    data.write(34, 5, 3, table_format)
    data.write(35, 5, 2, table_format)

    # likelihood and risk sheet
    likelihoodAndRisk.merge_range(2, 2, 2, 7,
                                  "(U) ARL/SLAD employs the Cybersecurity Risk Assessment Guide v1.1 Standard Operation Procedure for the ARL/SLAD Information and Electronic Protection Division (IEPD). Determining the Assessed Risk requires several inputs:",
                                  table_format)
    likelihoodAndRisk.set_column(2, 2, float(30))
    likelihoodAndRisk.write(3, 2, "1. (U) Vulnerability Severity (Vs)", table_format)
    likelihoodAndRisk.write(4, 2, "2. (U) Sraw", table_format)
    likelihoodAndRisk.write(5, 2, "3. (U) Impact", table_format)
    likelihoodAndRisk.write(6, 2, "4. (U) Countermeasure", table_format)
    likelihoodAndRisk.write(7, 2, "5. (U) Relevance of Threat", table_format)
    likelihoodAndRisk.write(8, 2, "6. (U) Likelihood", table_format)
    likelihoodAndRisk.merge_range(10, 2, 10, 7,
                                  " (U) The Vulnerability Severity (Vs) measures the potential impact a vulnerability has on a system and is calculated by the following formula in Figure X0:")
    likelihoodAndRisk.write(11, 2, "Vs = (Sraw * Impact * Countermeasure) / 10")
    likelihoodAndRisk.write(12, 2, "UNCLASSIFIED", table_format)
    likelihoodAndRisk.merge_range(14, 2, 14, 7,
                                  "(U) Sraw is directly obtained from the Security Technical Implementation Guide (STIG)/Security Reference Guide (SRG) or assigned by the Security Control Assessor Validator (SCA-V). Discovered vulnerabilities are assigned qualitative DISA CAT codes as described in the DISA STIG on Application Security Development, shown in Table X1.",
                                  table_format)

    # relevance of threat sheet
    relevanceOfThreat.merge_range(1, 1, 1, 2, "Table X7 (U) Relevance of Threat", table_format)
    relevanceOfThreat.set_column(1, 1, float(8))
    relevanceOfThreat.set_column(1, 2, float(30))
    relevanceOfThreat.write(2, 1, "Confirmed", table_format)
    relevanceOfThreat.write(2, 2, "The threat event or TTP has been seen by the organization.", table_format)
    relevanceOfThreat.write(3, 1, "Expected", table_format)
    relevanceOfThreat.write(3, 2, "The threat event or TTP has been seen by the organization’s peers or partners.",
                            table_format)
    relevanceOfThreat.write(4, 1, "Anticipated", table_format)
    relevanceOfThreat.write(4, 2, "The threat event or TTP has been reported by a trusted source.", table_format)
    relevanceOfThreat.write(5, 1, "Predicted", table_format)
    relevanceOfThreat.write(5, 2, "The threat event or TTP has been predicted by a trusted source.", table_format)
    relevanceOfThreat.write(6, 1, "Possible", table_format)
    relevanceOfThreat.write(6, 2, "The threat event or TTP has been described by a somewhat credible source.",
                            table_format)
    relevanceOfThreat.merge_range(7, 1, 7, 2, "UNCLASSIFIED", table_format)

    relevanceOfThreat.merge_range(9, 1, 14, 2,
                                  "(U) For this assessment, ARL/SLAD relied on analyst experiences and observations while supporting Adversarial Assessments (AAs) to derive the Relevance of Threat value. ",
                                  table_format)

    workbook.close()
